"""
AI Expert System - PPT Parser Module
解析 PowerPoint 檔案，提取文字與內嵌圖片
"""

import os
import logging
from typing import List, Dict
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

logger = logging.getLogger(__name__)


def extract_text_from_slide(slide) -> str:
    """提取單張投影片的所有文字"""
    text_runs = []
    
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            text_runs.append(shape.text.strip())
    
    return "\n".join(text_runs)


def extract_embedded_images(slide, output_dir: str, slide_id: int) -> List[str]:
    """
    提取投影片中的內嵌圖片（不做整頁截圖）
    僅當圖片已存在於 PPT 中時才提取
    
    Args:
        slide: python-pptx slide object
        output_dir: 圖片儲存目錄
        slide_id: 投影片 ID (用於命名)
    
    Returns:
        List[str]: 提取的圖片路徑列表
    """
    os.makedirs(output_dir, exist_ok=True)
    image_paths = []
    
    for idx, shape in enumerate(slide.shapes):
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_bytes = image.blob
                image_ext = image.ext
                
                # 儲存圖片
                image_filename = f"slide_{slide_id}_img_{idx}.{image_ext}"
                image_path = os.path.join(output_dir, image_filename)
                
                with open(image_path, 'wb') as f:
                    f.write(image_bytes)
                
                image_paths.append(image_path)
                logger.debug(f"提取圖片: {image_filename}")
        except Exception as e:
            logger.warning(f"提取圖片失敗 (投影片 {slide_id}, shape {idx}): {e}")
            continue
    
    return image_paths


def parse_ppt(file_path: str, extract_images: bool = True) -> List[Dict]:
    """
    解析 PPT 檔案
    
    Args:
        file_path: PPT 檔案路徑
        extract_images: 是否提取內嵌圖片
    
    Returns:
        List[Dict]: 投影片資料列表
        格式: [
            {
                'page_num': 1,
                'text': '投影片文字內容',
                'images': ['path/to/img1.png', ...]  # 若 extract_images=False 則為空
            },
            ...
        ]
    """
    logger.info(f"{'[文字+圖片]' if extract_images else '[僅文字]'} 開始解析: {file_path}")
    
    try:
        prs = Presentation(file_path)
        slides_data = []
        
        output_dir = "data/temp_images"
        
        for idx, slide in enumerate(prs.slides):
            page_num = idx + 1
            
            # 提取文字
            text = extract_text_from_slide(slide)
            
            # 提取圖片（若啟用）
            images = []
            if extract_images:
                try:
                    images = extract_embedded_images(slide, output_dir, slide.slide_id)
                except Exception as e:
                    logger.warning(f"提取圖片失敗 (頁 {page_num}): {e}")
            
            slides_data.append({
                'page_num': page_num,
                'text': text,
                'images': images
            })
            
            logger.debug(f"頁 {page_num}: {len(text)} 字元, {len(images)} 圖片")
        
        logger.info(f"✅ 解析完成: {len(slides_data)} 頁, 檔案: {os.path.basename(file_path)}")
        return slides_data
        
    except Exception as e:
        logger.error(f"❌ 解析失敗: {file_path}, 錯誤: {e}")
        raise


# 測試用
if __name__ == "__main__":
    import sys
    
    if len(sys.argv) > 1:
        test_file = sys.argv[1]
        result = parse_ppt(test_file, extract_images=True)
        print(f"解析結果: {len(result)} 頁")
        for slide in result[:3]:  # 顯示前3頁
            print(f"\n頁 {slide['page_num']}:")
            print(f"  文字: {slide['text'][:100]}...")
            print(f"  圖片: {len(slide['images'])} 張")
    else:
        print("使用方式: python ppt_parser.py <ppt_file_path>")
